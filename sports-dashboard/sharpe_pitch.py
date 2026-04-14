#!/usr/bin/env python3
"""Generate the Sharpe pitch deck — minimalist black & white aesthetic."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xA3, 0xA3, 0xA3)
LIGHT_GRAY = RGBColor(0xE5, 0xE5, 0xE5)
OFF_WHITE = RGBColor(0xFA, 0xFA, 0xFA)
GREEN = RGBColor(0x16, 0xA3, 0x4A)

FONT = "Inter"
FONT_LIGHT = "Inter"

W = Inches(13.333)
H = Inches(7.5)


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=14, color=BLACK, bold=False, align=PP_ALIGN.LEFT, font=FONT, anchor=MSO_ANCHOR.TOP, spacing=None, line_spacing=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    if spacing:
        p.space_after = Pt(spacing)
    return txBox


def add_multiline(slide, left, top, width, height, lines, size=14, color=BLACK, bold=False, align=PP_ALIGN.LEFT, line_spacing=None):
    """lines is list of (text, size, color, bold) tuples"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, line_data in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(line_data, str):
            p.text = line_data
            p.font.size = Pt(size)
            p.font.color.rgb = color
            p.font.bold = bold
        else:
            p.text = line_data[0]
            p.font.size = Pt(line_data[1]) if len(line_data) > 1 else Pt(size)
            p.font.color.rgb = line_data[2] if len(line_data) > 2 else color
            p.font.bold = line_data[3] if len(line_data) > 3 else bold
        p.font.name = FONT
        p.alignment = align
        if line_spacing:
            p.line_spacing = Pt(line_spacing)
    return txBox


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.fill.solid()
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    return shape


def add_line(slide, x1, y1, x2, y2, color=LIGHT_GRAY, width=1):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = straight
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector


# ============================================================
# SLIDE 1 — Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide, BLACK)

add_text(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(1.2),
         "SHARPE", size=72, color=WHITE, bold=False, align=PP_ALIGN.LEFT)

add_text(slide, Inches(1.5), Inches(3.5), Inches(8), Inches(0.6),
         "Sports Market Intelligence", size=24, color=GRAY, bold=False)

add_text(slide, Inches(1.5), Inches(4.4), Inches(8), Inches(0.8),
         "Find mispriced bets before the market corrects.\nReal-time comparison of bookmaker odds, Polymarket, and Kalshi.",
         size=14, color=GRAY, line_spacing=22)

# Thin horizontal line
add_line(slide, Inches(1.5), Inches(5.8), Inches(5), Inches(5.8), color=GRAY, width=0.5)

add_text(slide, Inches(1.5), Inches(6.0), Inches(4), Inches(0.4),
         "2026", size=12, color=GRAY)


# ============================================================
# SLIDE 2 — Problem
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, OFF_WHITE)

add_text(slide, Inches(1.5), Inches(0.8), Inches(3), Inches(0.4),
         "THE PROBLEM", size=11, color=GRAY, bold=False)

add_text(slide, Inches(1.5), Inches(1.5), Inches(9), Inches(1.5),
         "Billions left on the table.\nBookmaker odds and prediction\nmarkets rarely talk to each other.",
         size=40, color=BLACK, bold=False, line_spacing=52)

add_line(slide, Inches(1.5), Inches(4.2), Inches(11.8), Inches(4.2), color=LIGHT_GRAY)

# Three problem cards
problems = [
    ("$100B+", "Global online gambling market\nwith systemic pricing inefficiencies"),
    ("Fragmented", "Existing tools are basic converters,\nnot real-time intelligence"),
    ("4-10%", "Bookmaker vig creates structural\npricing gaps vs. prediction markets"),
]

for i, (stat, desc) in enumerate(problems):
    x = Inches(1.5 + i * 3.5)
    add_text(slide, x, Inches(4.6), Inches(3), Inches(0.8),
             stat, size=44, color=BLACK, bold=False)
    add_text(slide, x, Inches(5.5), Inches(3), Inches(0.8),
             desc, size=13, color=GRAY, line_spacing=20)


# ============================================================
# SLIDE 3 — Solution
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, OFF_WHITE)

add_text(slide, Inches(1.5), Inches(0.8), Inches(3), Inches(0.4),
         "THE SOLUTION", size=11, color=GRAY)

add_text(slide, Inches(1.5), Inches(1.5), Inches(9), Inches(1.2),
         "One dashboard.\nThree price sources.\nEvery edge visible.",
         size=40, color=BLACK, line_spacing=52)

add_line(slide, Inches(1.5), Inches(4.0), Inches(11.8), Inches(4.0), color=LIGHT_GRAY)

add_text(slide, Inches(1.5), Inches(4.4), Inches(10), Inches(0.8),
         "Sharpe scans bookmaker odds, Polymarket, and Kalshi in real time to surface\npricing divergences and arbitrage opportunities across 25+ sports and esports.",
         size=15, color=GRAY, line_spacing=24)

# Three source boxes
sources = [
    ("BOOKMAKERS", "Aggregated odds from\ntop sportsbooks worldwide"),
    ("POLYMARKET", "Decentralized prediction\nmarket pricing"),
    ("KALSHI", "1,100+ regulated\nevent contracts"),
]

for i, (title, desc) in enumerate(sources):
    x = Inches(1.5 + i * 3.5)
    add_rect(slide, x, Inches(5.5), Inches(3), Inches(1.4), fill_color=WHITE, line_color=LIGHT_GRAY)
    add_text(slide, x + Inches(0.3), Inches(5.7), Inches(2.4), Inches(0.3),
             title, size=11, color=BLACK, bold=True)
    add_text(slide, x + Inches(0.3), Inches(6.1), Inches(2.4), Inches(0.6),
             desc, size=12, color=GRAY, line_spacing=18)


# ============================================================
# SLIDE 4 — How It Works
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, OFF_WHITE)

add_text(slide, Inches(1.5), Inches(0.8), Inches(3), Inches(0.4),
         "HOW IT WORKS", size=11, color=GRAY)

add_text(slide, Inches(1.5), Inches(1.5), Inches(9), Inches(0.8),
         "From data to edge in seconds.", size=40, color=BLACK)

steps = [
    ("01", "AGGREGATE", "Scan odds from top bookmakers\nworldwide via The Odds API.\nCover 25+ sports leagues."),
    ("02", "COMPARE", "Pull real-time prices from\nPolymarket and Kalshi.\n1,100+ active markets."),
    ("03", "ANALYZE", "Kelly Criterion sizing and\ndivergence analysis surface\nmispriced opportunities."),
    ("04", "ACT", "Confidence scoring, alerts,\nand direct links to trade\non each platform."),
]

for i, (num, title, desc) in enumerate(steps):
    x = Inches(1.5 + i * 2.8)
    y = Inches(3.0)

    # Number
    add_rect(slide, x, y, Inches(0.5), Inches(0.5), fill_color=BLACK)
    add_text(slide, x, y, Inches(0.5), Inches(0.5),
             num, size=14, color=WHITE, bold=False, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(slide, x, y + Inches(0.8), Inches(2.4), Inches(0.3),
             title, size=11, color=BLACK, bold=True)
    add_text(slide, x, y + Inches(1.2), Inches(2.4), Inches(1.0),
             desc, size=12, color=GRAY, line_spacing=18)

    # Arrow between steps
    if i < 3:
        arrow_x = x + Inches(2.55)
        add_text(slide, arrow_x, y, Inches(0.3), Inches(0.5),
                 "\u2192", size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ============================================================
# SLIDE 5 — Coverage
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, OFF_WHITE)

add_text(slide, Inches(1.5), Inches(0.8), Inches(3), Inches(0.4),
         "COVERAGE", size=11, color=GRAY)

add_text(slide, Inches(1.5), Inches(1.5), Inches(9), Inches(0.8),
         "25+ leagues. Every major market.", size=40, color=BLACK)

# Sports grid
sports_col1 = [
    "NBA", "NFL", "NHL", "MLB", "NCAAF",
    "MLS", "Boxing"
]
sports_col2 = [
    "EPL", "La Liga", "Bundesliga",
    "Serie A", "Ligue 1",
    "Champions League", "Europa League"
]
sports_col3 = [
    "Tennis ATP", "Tennis WTA",
    "Formula 1", "LoL (LCK/LPL/LEC)",
    "CS2", "Valorant", "Dota 2"
]

headers = ["US SPORTS", "EUROPEAN FOOTBALL", "OTHER / ESPORTS"]
cols = [sports_col1, sports_col2, sports_col3]

for i, (header, items) in enumerate(zip(headers, cols)):
    x = Inches(1.5 + i * 3.5)

    add_text(slide, x, Inches(3.0), Inches(3), Inches(0.3),
             header, size=10, color=GRAY, bold=False)

    add_line(slide, x, Inches(3.4), x + Inches(3), Inches(3.4), color=LIGHT_GRAY)

    for j, sport in enumerate(items):
        add_text(slide, x, Inches(3.6 + j * 0.42), Inches(3), Inches(0.35),
                 sport, size=14, color=BLACK)


# ============================================================
# SLIDE 6 — Product Features
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, OFF_WHITE)

add_text(slide, Inches(1.5), Inches(0.8), Inches(3), Inches(0.4),
         "PRODUCT", size=11, color=GRAY)

add_text(slide, Inches(1.5), Inches(1.5), Inches(9), Inches(0.8),
         "Built for serious bettors.", size=40, color=BLACK)

features = [
    ("THREE-SOURCE COMPARISON", "Real-time prices from bookmakers,\nPolymarket, and Kalshi side by side"),
    ("DIVERGENCE ALERTS", "Confidence-scored signals when\npricing gaps exceed your threshold"),
    ("PLAYER PROPS", "NBA points, assists, rebounds,\n3-pointers from Kalshi"),
    ("FUTURES & AWARDS", "MVP, ROY, Heisman, league\nwinners across all sports"),
    ("PROFIT TRACKER", "Log trades, track P&L, measure\nyour edge over time"),
    ("WATCHLIST", "Save markets and get notified\nwhen edges appear"),
]

for i, (title, desc) in enumerate(features):
    col = i % 3
    row = i // 3
    x = Inches(1.5 + col * 3.5)
    y = Inches(3.0 + row * 2.0)

    add_rect(slide, x, y, Inches(3), Inches(1.6), fill_color=WHITE, line_color=LIGHT_GRAY)
    add_text(slide, x + Inches(0.3), y + Inches(0.25), Inches(2.4), Inches(0.3),
             title, size=10, color=BLACK, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(0.65), Inches(2.4), Inches(0.7),
             desc, size=12, color=GRAY, line_spacing=18)


# ============================================================
# SLIDE 7 — Market Opportunity
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BLACK)

add_text(slide, Inches(1.5), Inches(0.8), Inches(3), Inches(0.4),
         "MARKET", size=11, color=GRAY)

add_text(slide, Inches(1.5), Inches(1.5), Inches(9), Inches(1.2),
         "Two massive markets.\nNo real-time intelligence layer.",
         size=40, color=WHITE, line_spacing=52)

# Market stats
stats = [
    ("$100B+", "Global online\ngambling market"),
    ("$25B+", "Polymarket cumulative\ntrading volume"),
    ("$24B+", "Kalshi 2025\ntrading volume"),
    ("None", "Purpose-built real-time\nintelligence dashboards"),
]

for i, (val, label) in enumerate(stats):
    x = Inches(1.5 + i * 2.8)
    y = Inches(4.0)

    add_text(slide, x, y, Inches(2.4), Inches(0.8),
             val, size=48, color=WHITE, bold=False)
    add_text(slide, x, y + Inches(0.9), Inches(2.4), Inches(0.6),
             label, size=13, color=GRAY, line_spacing=20)

    if i < 3:
        add_line(slide, x + Inches(2.6), y, x + Inches(2.6), y + Inches(1.5),
                 color=RGBColor(0x40, 0x40, 0x40))


# ============================================================
# SLIDE 8 — Business Model
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, OFF_WHITE)

add_text(slide, Inches(1.5), Inches(0.8), Inches(3), Inches(0.4),
         "BUSINESS MODEL", size=11, color=GRAY)

add_text(slide, Inches(1.5), Inches(1.5), Inches(9), Inches(0.8),
         "Freemium with clear upgrade path.", size=40, color=BLACK)

# Free tier
add_rect(slide, Inches(1.5), Inches(3.0), Inches(5), Inches(3.8), fill_color=WHITE, line_color=LIGHT_GRAY)
add_text(slide, Inches(1.8), Inches(3.2), Inches(4), Inches(0.4),
         "FREE", size=11, color=GRAY, bold=False)
add_text(slide, Inches(1.8), Inches(3.6), Inches(2), Inches(0.6),
         "$0", size=36, color=BLACK, bold=False)
add_text(slide, Inches(3.3), Inches(3.8), Inches(2), Inches(0.4),
         "forever", size=14, color=GRAY)

free_features = [
    "4 major US sports",
    "Basic divergence data",
    "3 cards visible per sport",
    "Community access",
]
for j, feat in enumerate(free_features):
    add_text(slide, Inches(1.8), Inches(4.5 + j * 0.42), Inches(4), Inches(0.35),
             feat, size=13, color=GRAY)

# Pro tier
add_rect(slide, Inches(7), Inches(3.0), Inches(5), Inches(3.8), fill_color=BLACK)
add_text(slide, Inches(7.3), Inches(3.2), Inches(4), Inches(0.4),
         "PRO", size=11, color=GRAY, bold=False)
add_text(slide, Inches(7.3), Inches(3.6), Inches(2), Inches(0.6),
         "$19", size=36, color=WHITE, bold=False)
add_text(slide, Inches(8.8), Inches(3.8), Inches(2), Inches(0.4),
         "/month", size=14, color=GRAY)

pro_features = [
    "25+ sports and esports",
    "All data points and intel",
    "Unlimited market cards",
    "Divergence alerts",
    "Profit tracker",
    "Priority support",
]
for j, feat in enumerate(pro_features):
    add_text(slide, Inches(7.3), Inches(4.5 + j * 0.42), Inches(4), Inches(0.35),
             feat, size=13, color=RGBColor(0xCC, 0xCC, 0xCC))


# ============================================================
# SLIDE 9 — Technology
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, OFF_WHITE)

add_text(slide, Inches(1.5), Inches(0.8), Inches(3), Inches(0.4),
         "TECHNOLOGY", size=11, color=GRAY)

add_text(slide, Inches(1.5), Inches(1.5), Inches(9), Inches(0.8),
         "Lean architecture. Fast iteration.", size=40, color=BLACK)

tech_items = [
    ("BACKEND", "Single-file Python (FastAPI + Uvicorn)\nReal-time WebSocket push to all clients\nAsync parallel API fetching"),
    ("DATA SOURCES", "The Odds API — bookmaker odds\nPolymarket Gamma API — prediction markets\nKalshi Trade API — 34 series, 1,100+ markets"),
    ("INFRASTRUCTURE", "Ubuntu server deployment\nCloudflare Tunnels for HTTPS\nSQLite user database"),
    ("FRONTEND", "Vanilla JS — zero dependencies\nMinimalist responsive design\nLive updates, no page refreshes"),
]

for i, (title, desc) in enumerate(tech_items):
    col = i % 2
    row = i // 2
    x = Inches(1.5 + col * 5.5)
    y = Inches(3.0 + row * 2.0)

    add_text(slide, x, y, Inches(4.5), Inches(0.3),
             title, size=10, color=GRAY, bold=False)
    add_line(slide, x, y + Inches(0.35), x + Inches(4.5), y + Inches(0.35), color=LIGHT_GRAY)
    add_text(slide, x, y + Inches(0.5), Inches(4.5), Inches(1.0),
             desc, size=13, color=BLACK, line_spacing=20)


# ============================================================
# SLIDE 10 — Ask / Close
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BLACK)

add_text(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
         "SHARPE", size=72, color=WHITE, bold=False)

add_text(slide, Inches(1.5), Inches(3.5), Inches(8), Inches(0.6),
         "The edge between odds and markets.", size=24, color=GRAY)

add_line(slide, Inches(1.5), Inches(4.8), Inches(5), Inches(4.8), color=GRAY, width=0.5)

add_text(slide, Inches(1.5), Inches(5.2), Inches(8), Inches(0.8),
         "Let's talk.\nsharpe.app", size=16, color=GRAY, line_spacing=26)


# Save
out = Path(__file__).parent / "Sharpe_Pitch.pptx"
prs.save(str(out))
print(f"Saved: {out}")
